#!/usr/bin/env python3
# Copyright (c) OpenAI. All rights reserved.
import argparse
import sys
import tempfile
from os.path import abspath, expanduser, join
from pathlib import Path
from typing import Sequence, cast

import numpy as np

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

import render_slides  # type: ignore
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu

# Configuration specific to overflow checking
PAD_PX: int = 100  # fixed padding on every side in pixels
PAD_RGB = (200, 200, 200)
EMU_PER_INCH: int = 914_400


def px_to_emu(px: int, dpi: int) -> Emu:
    return Emu(int(px * EMU_PER_INCH // dpi))


def calc_tol(dpi: int) -> int:
    """Calculate per-channel colour tolerance appropriate for *dpi* (anti-aliasing tolerance)."""
    if dpi >= 300:
        return 0
    # 1 at 250 DPI, 5 at 150 DPI, capped to 10.
    tol = round((300 - dpi) / 25)
    return min(max(tol, 1), 10)


def enlarge_deck(src: str, dst: str, pad_emu: Emu) -> tuple[int, int]:
    """Enlarge the input PPTX with a fixed grey padding and return the new page size."""
    prs = Presentation(src)
    w0 = cast(Emu, prs.slide_width)
    h0 = cast(Emu, prs.slide_height)
    w1 = Emu(w0 + 2 * pad_emu)
    h1 = Emu(h0 + 2 * pad_emu)
    prs.slide_width = w1
    prs.slide_height = h1

    for slide in prs.slides:
        # Shift all shapes so the original canvas sits centred in the new deck.
        for shp in list(slide.shapes):
            shp.left = Emu(int(shp.left) + pad_emu)
            shp.top = Emu(int(shp.top) + pad_emu)

        pads = (
            (Emu(0), Emu(0), pad_emu, h1),  # left
            (Emu(int(w1) - int(pad_emu)), Emu(0), pad_emu, h1),  # right
            (Emu(0), Emu(0), w1, pad_emu),  # top
            (Emu(0), Emu(int(h1) - int(pad_emu)), w1, pad_emu),  # bottom
        )

        sp_tree = slide.shapes._spTree  # pylint: disable=protected-access

        for left, top, width, height in pads:
            pad_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            pad_shape.fill.solid()
            pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
            pad_shape.line.fill.background()

            # Send pad behind all other shapes (index 2 after mandatory nodes)
            sp_tree.remove(pad_shape._element)
            sp_tree.insert(2, pad_shape._element)

    prs.save(dst)
    return int(w1), int(h1)


def inspect_images(
    paths: Sequence[str],
    pad_ratio_w: float,
    pad_ratio_h: float,
    dpi: int,
) -> list[int]:
    """Return 1-based indices of slides that contain pixels outside the pad."""

    tol = calc_tol(dpi)
    failures: list[int] = []
    pad_colour = np.array(PAD_RGB, dtype=np.uint8)

    for idx, img_path in enumerate(paths, start=1):
        with Image.open(img_path) as img:
            rgb = img.convert("RGB")
            arr = np.asarray(rgb)

        h, w, _ = arr.shape
        # Exclude the innermost 1-pixel band
        pad_x = int(w * pad_ratio_w) - 1
        pad_y = int(h * pad_ratio_h) - 1

        left_margin = arr[:, :pad_x, :]
        right_margin = arr[:, w - pad_x :, :]
        top_margin = arr[:pad_y, :, :]
        bottom_margin = arr[h - pad_y :, :, :]

        def _is_clean(margin: np.ndarray) -> bool:
            diff = np.abs(margin.astype(np.int16) - pad_colour)
            matches = np.all(diff <= tol, axis=-1)
            mismatch_fraction = 1.0 - (np.count_nonzero(matches) / matches.size)
            if dpi >= 300:
                max_mismatch = 0.01
            elif dpi >= 200:
                max_mismatch = 0.02
            else:
                max_mismatch = 0.03
            return mismatch_fraction <= max_mismatch

        if not (
            _is_clean(left_margin)
            and _is_clean(right_margin)
            and _is_clean(top_margin)
            and _is_clean(bottom_margin)
        ):
            failures.append(idx)

    return failures


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Check a PPTX for content overflowing the original canvas by rendering with padding "
            "and inspecting the margins."
        )
    )
    parser.add_argument(
        "input_path",
        type=str,
        help="Path to the input PPTX file.",
    )
    parser.add_argument(
        "--width",
        type=int,
        default=1600,
        help=(
            "Approximate maximum width in pixels after isotropic scaling (default 1600). "
            "The actual value may exceed slightly."
        ),
    )
    parser.add_argument(
        "--height",
        type=int,
        default=900,
        help=(
            "Approximate maximum height in pixels after isotropic scaling (default 900). "
            "The actual value may exceed slightly."
        ),
    )
    parser.add_argument(
        "--pad_px",
        type=int,
        default=PAD_PX,
        help="Padding in pixels to add on each side before rasterization.",
    )
    args = parser.parse_args()

    input_path = abspath(expanduser(args.input_path))
    # Width and height refer to the original, unaltered slide dimensions.
    dpi = render_slides.calc_dpi_via_ooxml(input_path, args.width, args.height)

    # Not using ``tempfile.TemporaryDirectory(delete=False)`` for Python 3.11 compatibility.
    tmpdir = tempfile.mkdtemp()
    enlarged_pptx = join(tmpdir, "enlarged.pptx")
    pad_emu = px_to_emu(args.pad_px, dpi)
    w1, h1 = enlarge_deck(input_path, enlarged_pptx, pad_emu=pad_emu)
    pad_ratio_w = pad_emu / w1
    pad_ratio_h = pad_emu / h1

    img_dir = join(tmpdir, "imgs")
    img_paths = render_slides.rasterize(enlarged_pptx, img_dir, dpi)
    failing = inspect_images(img_paths, pad_ratio_w, pad_ratio_h, dpi)

    if failing:
        print(
            "ERROR: Slides with content overflowing original canvas (1-based indexing): "
            + ", ".join(map(str, failing))
            + "\n"
            + "Rendered images with grey paddings for problematic slides are available at: "
        )
        for i in failing:
            print(img_paths[i - 1])
    else:
        print("Test passed. No overflow detected.")


if __name__ == "__main__":
    main()
